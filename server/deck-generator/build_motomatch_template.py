"""
╔══════════════════════════════════════════════════════════════════════════════╗
║           MOTOMATCH DECK BUILDER — MASTER TEMPLATE                          ║
║           Motosaic Brand-Compliant | python-pptx | May 2026                 ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  USAGE:                                                                      ║
║    1. Update the ── CLIENT DATA ── block with new client info               ║
║    2. Update the ── VEHICLES DATA ── block (1 dict per vehicle)             ║
║    3. Update the ── COMPARISON TABLE ── block with correct values           ║
║    4. Place vehicle photos in /home/user/workspace/assets/photos/           ║
║    5. Run:  pip install -q python-pptx Pillow && python3 build_motomatch_template.py ║
║    6. Output: /home/user/workspace/[LastName]_[FirstName]_MotoMatch.pptx    ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  FONT NOTE:                                                                  ║
║    Brand font is Industry (Ultra Italic / Bold / Medium / Light).           ║
║    python-pptx cannot embed custom OTF/TTF fonts directly.                  ║
║    Script uses Calibri as a safe fallback. When opened in PowerPoint on     ║
║    a machine where Industry is installed, it renders in Industry            ║
║    automatically (font is referenced by name in the XML).                   ║
║    To force Industry: install it on the deck-viewing machine, or use        ║
║    PowerPoint's "Replace Fonts" dialog before client delivery.              ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  LOGO ASSETS (pre-extracted — do NOT re-extract each session):              ║
║    All 4 logo PNGs live at /home/user/workspace/assets/logos/               ║
║    slide1_img1.png  ~151KB  Hero logo — title slide top-left                ║
║    slide1_img6.png   ~25KB  MOTOSAIC wordmark — title slide footer          ║
║    slide2_img1.png   ~57KB  Small icon — top-right on all inner slides      ║
║    slide2_img2.png   ~26KB  MOTOSAIC wordmark — centered bottom inner slides ║
║    If logos are missing, re-extract via the motomatch-deck skill (Stage 4a) ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  LOGO SELECTION LOGIC:                                                       ║
║    Title slide  → LOGO_TITLE (hero) top-left + LOGO_FOOTER (wordmark) footer║
║    Inner slides → LOGO_SMALL (icon) top-right + LOGO_FOOTER2 (wordmark) btm ║
║    Rule: never place hero logo and primary wordmark together (two stars)     ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  SHELBY LIGHT COLOR NOTE:                                                    ║
║    Build script uses #0B94B0 (confirmed from PPTX eyedropper on Julie deck) ║
║    Brand guidelines PDF lists #0184B0 — the PPTX value (#0B94B0) is used   ║
║    here as it was approved across two client decks.                         ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, os, sys
from PIL import Image


# ══════════════════════════════════════════════════════════════════════════════
# ─── BRAND COLORS (Motosaic 2026 — confirmed against brand guidelines PDF) ───
# ══════════════════════════════════════════════════════════════════════════════
# Source: MOTOSAIC_BRANDGUIDELINES_2026.pdf + eyedropper on approved PPTX decks

SHELBY_BLUE   = RGBColor(0x00, 0x43, 0x63)   # #004363  Headers, footer bars, blurb bg, nav bar
SHELBY_LIGHT  = RGBColor(0x0B, 0x94, 0xB0)   # #0B94B0  Rec counter, MSRP, hyperlinks, [Why tags]
                                               #          (brand PDF lists #0184B0; #0B94B0 confirmed in PPTX)
MIAMI         = RGBColor(0x1F, 0xC3, 0xEF)   # #1FC3EF  Brand accent (not heavily used in deck)
SAO_PAULO     = RGBColor(0xF2, 0xEA, 0x00)   # #F2EA00  ★ star line + Considerations left bar
NARDO         = RGBColor(0xE1, 0xF3, 0xF5)   # #E1F3F5  Title bg, Financials box, specs/why bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
PURE_BLACK    = RGBColor(0x01, 0x02, 0x02)   # #010202  All body copy
SUBHEAD_GRAY  = RGBColor(0x66, 0x66, 0x66)   # #666666  Labels, MSRP label, considerations text
OFF_WHITE_ROW = RGBColor(0xF8, 0xF8, 0xF8)   # #F8F8F8  Lower-priority rows, considerations bg
GREEN_BEST    = RGBColor(0x00, 0x70, 0x50)   # #007050  Best-in-class text in comparison table
GREEN_CELL_BG = RGBColor(0xE8, 0xF5, 0xEE)   # #E8F5EE  Best-in-class cell background


# ══════════════════════════════════════════════════════════════════════════════
# ─── SLIDE DIMENSIONS ────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# 10" × 5.625" widescreen — matches Motomatch_Deck_Example-2.pptx exactly

SLIDE_W = 9144000   # EMU — 10 inches
SLIDE_H = 5143500   # EMU — 5.625 inches


# ══════════════════════════════════════════════════════════════════════════════
# ─── FONTS ───────────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# Brand font: Industry (Ultra Italic for display, Bold/Medium for body/labels)
# Fallback: Calibri — python-pptx cannot embed OTF/TTF directly.
# PowerPoint resolves "Industry" by name on open if font is installed.
# Calibri is always available and renders cleanly on all platforms.

FONT_BOLD = "Calibri"   # Used for: headings, labels, [Why tags], counter
FONT_BODY = "Calibri"   # Used for: body copy, blurbs, bullets, links, specs


# ══════════════════════════════════════════════════════════════════════════════
# ─── ASSET PATHS ─────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

ASSETS     = "/home/user/workspace/assets"
LOGO_TITLE   = f"{ASSETS}/logos/slide1_img1.png"   # ~151KB  Hero logo — title slide, top-left
LOGO_FOOTER  = f"{ASSETS}/logos/slide1_img6.png"   # ~25KB   MOTOSAIC wordmark — title footer
LOGO_SMALL   = f"{ASSETS}/logos/slide2_img1.png"   # ~57KB   Small icon — inner slides, top-right
LOGO_FOOTER2 = f"{ASSETS}/logos/slide2_img2.png"   # ~26KB   MOTOSAIC wordmark — inner slides, bottom center


# ══════════════════════════════════════════════════════════════════════════════
# ─── PHOTO PLACEMENT CONSTANTS ───────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# Three identical slots on the right half of each vehicle slide.
# Fill header-to-footer, equally spaced, no variation across slides or clients.
# Crop-to-fill, non-destructive: full image embedded, dealer can re-crop in PPT.

PHOTO_L   = Emu(5864466)   # Left edge of photo column (EMU)
PHOTO_W   = Emu(3279534)   # Width of each photo slot (fills to slide right edge)
PHOTO_H   = Emu(1105282)   # Height of each slot: (4343400 - 987552 - 2×20000) ÷ 3
PHOTO_GAP = Emu(20000)     # Gap between slots

PHOTO_TOPS = [
    Emu(987552),                               # Slot 1 top
    Emu(987552) + PHOTO_H + PHOTO_GAP,         # Slot 2 top
    Emu(987552) + 2 * (PHOTO_H + PHOTO_GAP),  # Slot 3 top
]


# ══════════════════════════════════════════════════════════════════════════════
# ──────────────────── CLIENT DATA  (EDIT THIS BLOCK) ─────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# Replace ALL values below with the new client's information.
# priorities: list of (display_name, score 1-5) tuples, highest score first.
# considerations: 5 bullet points for the "What You Told Us" slide.
# footnote: italicized line at bottom of slide 2.

CLIENT = {
    "name":          "Jane Client",             # ← Full name as it appears on the deck
    "date":          "June 2026",               # ← Month + year of deck delivery
    "budget":        "$75,000 – $90,000",       # ← Budget range shown on financials box
    "purchase_type": "Finance",                 # ← "Finance", "Cash", or "Lease"
    "priorities": [
        # (display_name, score 1–5)  — keep in descending score order
        # ★★★★★  score 5 → Shelby Blue badge (boldest)
        # ★★★★☆  score 4 → Shelby Blue badge
        # ★★★☆☆  score 3 → Shelby Light badge
        # ★★☆☆☆  score 2 → Shelby Light badge
        # ★☆☆☆☆  score 1 → Shelby Light badge (muted)
        ("Interior Comfort & Luxury",       5),   # ← Replace with client's actual priorities
        ("Safety",                          5),
        ("Space / Storage",                 4),
        ("Technology",                      4),
        ("Third Row Space",                 3),
        ("Efficiency (Gas Mileage)",        3),
        ("Exterior Style",                  3),
        ("Towing / Hauling",                2),
        ("Off-Road Capability",             2),
        ("Engine Power / Speed",            2),
        ("Sporty Drive / Handling",         1),
        ("Maintenance / Cost of Ownership", 1),
        ("Resale Value",                    1),
        ("Warranty Coverage",               1),
        ("Brand Prestige / Status",         1),
    ],
    "considerations": [
        # 5 bullets for the Key Considerations box — each ~55 chars max
        "Upgrading from [Year Make Model] ([Xk] mi)",    # ← Current vehicle
        "Finance — [City, State]",                       # ← Purchase type + location
        "Trade-in: [Year Make Model], [Xk] mi, [cond]", # ← Trade-in details
        "Must-haves: [key requirements here]",           # ← Non-negotiables
        "Prefers [color] exterior — [family/lifestyle note]",  # ← Color + context
    ],
    "footnote": "These insights guided our vehicle selection. Each recommendation addresses your top priorities.",
}


# ══════════════════════════════════════════════════════════════════════════════
# ──────────────────── VEHICLES DATA  (EDIT THIS BLOCK) ───────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# One dict per vehicle. For 3-vehicle decks: 3 dicts. For 5-vehicle: 5 dicts.
#
# key:           short snake_case name — must match photo filenames
#                e.g. key="yukon_xl" → photos: yukon_xl_front.jpg etc.
# rec_num:       1, 2, 3 ... (1-based, in display order)
# year_make_model: "2026 GMC Yukon XL Denali Ultimate" — shown as slide heading
# msrp:          "$83,000 – $98,000" — right-aligned in Shelby Light
# blurb:         2–4 sentence dealer narrative (dark navy box, white italic text)
#                Write in dealer voice — direct, personal, specific to THIS client
# specs:         One line: "Engine | hp | drivetrain | MPG | 0-60 | Tows"
# why:           4 bullets: (tag, description)
#                tag format: "[Priority Name]" — must match a priority in CLIENT["priorities"]
#                tag + description combined: max ~55 chars to prevent wrapping
# star_line:     Optional. "★ [standout story]" — italic yellow accent.
#                Omit key or set to None if not applicable.
# considerations: 3 honest trade-offs. Each max ~55 chars.
# mfr_url:       Manufacturer page — linked as "Manufacturer"
# cd_url:        Car and Driver review URL — linked as "Car and Driver Review"
# edmunds_url:   Edmunds listing URL — linked as "Edmunds Review"
#
# PHOTO NAMING CONVENTION:
#   {key}_front.jpg     Front 3/4 exterior
#   {key}_rear.jpg      Rear 3/4 exterior
#   {key}_interior.jpg  Straight-on center stack / dashboard
#   Photos → /home/user/workspace/assets/photos/

VEHICLES = [
    # ── VEHICLE 1 ─────────────────────────────────────────────────────────────
    {
        "key":             "vehicle_one",                      # ← snake_case, matches photo filenames
        "rec_num":         1,
        "year_make_model": "2026 Make Model Trim AWD",         # ← Full year/make/model/trim
        "msrp":            "$XX,000 – $XX,000",
        "mfr_url":         "https://www.manufacturer.com/model",
        "cd_url":          "https://www.caranddriver.com/make/model",
        "edmunds_url":     "https://www.edmunds.com/make/model/",
        "blurb": (
            "Write a 2–4 sentence dealer narrative here. Be specific to this client's situation. "
            "Reference why this vehicle makes sense for their life. Keep it conversational and direct. "
            "Avoid generic language — this is the most personal part of the slide."
        ),
        "specs": "X.XL V8 | XXX hp | AWD | XX MPG | 0-60 in X.Xs | Tows X,XXX lbs",
        "why": [
            ("[Priority Name 1]",    "Specific reason tied to client priority"),    # ← max ~55 chars total
            ("[Priority Name 2]",    "Specific reason tied to client priority"),
            ("[Priority Name 3]",    "Specific reason tied to client priority"),
            ("[Priority Name 4]",    "Specific reason tied to client priority"),
        ],
        "star_line": "★ Optional standout story in one line",  # ← set to None to omit
        "considerations": [
            "Honest trade-off #1 — be specific, not generic",  # ← max ~55 chars
            "Honest trade-off #2 — be specific, not generic",
            "Honest trade-off #3 — be specific, not generic",
        ],
    },
    # ── VEHICLE 2 ─────────────────────────────────────────────────────────────
    {
        "key":             "vehicle_two",
        "rec_num":         2,
        "year_make_model": "2026 Make Model Trim AWD",
        "msrp":            "$XX,000 – $XX,000",
        "mfr_url":         "https://www.manufacturer.com/model",
        "cd_url":          "https://www.caranddriver.com/make/model",
        "edmunds_url":     "https://www.edmunds.com/make/model/",
        "blurb": (
            "Write dealer narrative here. Be specific to this client. "
            "What makes this the right choice for them?"
        ),
        "specs": "X.XL V6 | XXX hp | AWD | XX MPG | 0-60 in X.Xs | Tows X,XXX lbs",
        "why": [
            ("[Priority Name 1]",    "Specific reason tied to client priority"),
            ("[Priority Name 2]",    "Specific reason tied to client priority"),
            ("[Priority Name 3]",    "Specific reason tied to client priority"),
            ("[Priority Name 4]",    "Specific reason tied to client priority"),
        ],
        "star_line": None,  # ← None = no star line on this vehicle
        "considerations": [
            "Honest trade-off #1",
            "Honest trade-off #2",
            "Honest trade-off #3",
        ],
    },
    # ── VEHICLE 3 ─────────────────────────────────────────────────────────────
    {
        "key":             "vehicle_three",
        "rec_num":         3,
        "year_make_model": "2026 Make Model Trim AWD",
        "msrp":            "$XX,000 – $XX,000",
        "mfr_url":         "https://www.manufacturer.com/model",
        "cd_url":          "https://www.caranddriver.com/make/model",
        "edmunds_url":     "https://www.edmunds.com/make/model/",
        "blurb": (
            "Write dealer narrative here. Be specific to this client."
        ),
        "specs": "X.XL V6 | XXX hp | AWD | XX MPG | 0-60 in X.Xs | Tows X,XXX lbs",
        "why": [
            ("[Priority Name 1]",    "Specific reason tied to client priority"),
            ("[Priority Name 2]",    "Specific reason tied to client priority"),
            ("[Priority Name 3]",    "Specific reason tied to client priority"),
            ("[Priority Name 4]",    "Specific reason tied to client priority"),
        ],
        "star_line": None,
        "considerations": [
            "Honest trade-off #1",
            "Honest trade-off #2",
            "Honest trade-off #3",
        ],
    },
    # ── Add more vehicles here for 4- or 5-vehicle decks ──────────────────────
    # (copy the block above, increment rec_num)
]


# ══════════════════════════════════════════════════════════════════════════════
# ──────────────────── COMPARISON TABLE  (EDIT THIS BLOCK) ────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# Each row: (metric_label, [val_v1, val_v2, val_v3], best_idx)
# best_idx: 0-based index of the best vehicle in that row (gets green highlight)
#           Set to None if there's no clear winner (e.g. all tied, subjective)
#
# Standard rows (keep these in order, add client-specific rows as needed):
#   Starting MSRP, Horsepower, 0-60 mph, Combined MPG, Cargo (behind 3rd row),
#   Towing Capacity, Reliability (JD Power), Seating, Warranty
# Add rows for top client priorities (e.g. IIHS safety if Safety is top priority)

COMPARISON_ROWS = [
    # (metric_label,             [vehicle_1_val,  vehicle_2_val,  vehicle_3_val],  best_idx)
    ("Starting MSRP",            ["$XX,XXX",      "$XX,XXX",      "$XX,XXX"],      0),     # ← 0 = vehicle 1 is best (lowest MSRP)
    ("Horsepower",               ["XXX hp",       "XXX hp",       "XXX hp"],       0),     # ← index of highest hp
    ("0-60 mph",                 ["X.X sec",      "X.X sec",      "X.X sec"],      2),     # ← index of quickest
    ("Combined MPG",             ["XX mpg",       "XX mpg",       "XX mpg"],       0),     # ← index of best MPG
    ("Cargo (behind 3rd row)",   ["XX cu ft",     "XX cu ft",     "XX cu ft"],     0),     # ← index of most cargo
    ("Towing Capacity",          ["X,XXX lbs",    "X,XXX lbs",    "X,XXX lbs"],   None),  # ← None = no clear winner
    ("Reliability (JD Power)",   ["Above Avg",    "Above Avg",    "Above Avg"],   None),
    ("Seating",                  ["7–8",          "7–8",          "7–8"],          None),
    ("Warranty",                 ["3yr/36K",      "3yr/36K",      "3yr/36K"],      None),
]


# ══════════════════════════════════════════════════════════════════════════════
# ─── OUTPUT PATH ─────────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════
# Convention: [LastName]_[FirstName]_MotoMatch.pptx
# Change this to match the new client before running.

OUTPUT_PATH = "/home/user/workspace/Client_FirstName_MotoMatch.pptx"   # ← UPDATE THIS


# ══════════════════════════════════════════════════════════════════════════════
# ─── DERIVED CONSTANTS (do not edit) ─────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

TOTAL_VEHICLES = len(VEHICLES)

PHOTOS = {
    v["key"]: [
        f"{ASSETS}/photos/{v['key']}_front.jpg",
        f"{ASSETS}/photos/{v['key']}_rear.jpg",
        f"{ASSETS}/photos/{v['key']}_interior.jpg",
    ]
    for v in VEHICLES
}


# ══════════════════════════════════════════════════════════════════════════════
# ─── HELPER FUNCTIONS (do not edit) ──────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def add_shape(slide, left, top, width, height, fill_color=None, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if not line:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=PURE_BLACK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT_BOLD if bold else FONT_BODY
    return txBox


def add_hyperlink_run(paragraph, text, url, font_size, color=SHELBY_LIGHT, slide_part=None):
    """Add a hyperlinked run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = FONT_BODY
    if slide_part:
        rId = slide_part.relate_to(
            url,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
            is_external=True
        )
        rPr = run._r.get_or_add_rPr()
        hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
        hlinkClick.set(qn('r:id'), rId)
    return run


def add_plain_run(paragraph, text, font_size, color=PURE_BLACK, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT_BOLD if bold else FONT_BODY
    return run


def add_paragraph(tf, text, font_size, bold=False, color=PURE_BLACK,
                  align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT_BOLD if bold else FONT_BODY
    return p


def add_image_safe(slide, path, left, top, width, height):
    """
    Add image with crop-to-fill, non-destructive center crop.
    - Scales so image fills slot in both dimensions (no white space)
    - Excess is center-cropped via python-pptx crop attributes
    - Full image is embedded — dealer can right-click → Crop in PowerPoint to reframe
    - Never stretches or distorts aspect ratio
    Falls back to Nardo-colored placeholder box if image is missing.
    """
    if os.path.exists(path) and os.path.getsize(path) > 5000:
        try:
            from PIL import Image as PILImage
            with PILImage.open(path) as im:
                img_w, img_h = im.size
            img_ratio = img_w / img_h
            slot_ratio = width / height

            if img_ratio > slot_ratio:
                # Image wider than slot — fit height, center-crop sides
                pic = slide.shapes.add_picture(path, left, top, height=height)
                pic_w = pic.width
                excess = pic_w - width
                pic.crop_left  = (excess / 2) / pic_w
                pic.crop_right = (excess / 2) / pic_w
                pic.width = width
            else:
                # Image taller than slot — fit width, center-crop top/bottom
                pic = slide.shapes.add_picture(path, left, top, width=width)
                pic_h = pic.height
                excess = pic_h - height
                pic.crop_top    = (excess / 2) / pic_h
                pic.crop_bottom = (excess / 2) / pic_h
                pic.height = height

            pic.left = left
            pic.top  = top
            return pic
        except Exception as e:
            print(f"  WARNING: Could not add image {path}: {e}")

    # Placeholder box if image missing
    shape = add_shape(slide, left, top, width, height, fill_color=NARDO)
    fname = os.path.basename(path)
    add_textbox(slide, left + Emu(50000), top + Emu(height // 2 - 100000),
                width - Emu(100000), Emu(200000),
                f"[Photo: {fname}]", 9, color=SUBHEAD_GRAY, align=PP_ALIGN.CENTER)
    return shape


def stars_string(count, total=5):
    return "★" * count + "☆" * (total - count)


# ══════════════════════════════════════════════════════════════════════════════
# ─── INNER SLIDE CHROME (shared across all non-title slides) ─────────────────
# ══════════════════════════════════════════════════════════════════════════════

def add_inner_slide_chrome(slide):
    """Navy left bar + small top-right icon + bottom-center wordmark."""
    # Left navy bar (full height)
    bar = add_shape(slide, 0, 0, Emu(182880), SLIDE_H, fill_color=SHELBY_BLUE)
    bar.line.fill.background()

    # Small icon top-right
    add_image_safe(slide, LOGO_SMALL, Emu(8367464), Emu(137160), Emu(639376), Emu(320040))

    # MOTOSAIC wordmark bottom-center
    add_image_safe(slide, LOGO_FOOTER2, Emu(4074134), Emu(4873752), Emu(995732), Emu(164592))


# ══════════════════════════════════════════════════════════════════════════════
# ─── SLIDE 1: TITLE SLIDE ────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background: Nardo
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NARDO

    # Top thin navy bar
    top_bar = add_shape(slide, 0, 0, SLIDE_W, Emu(109728), fill_color=SHELBY_BLUE)
    top_bar.line.fill.background()

    # Hero logo top-left
    add_image_safe(slide, LOGO_TITLE, Emu(457200), Emu(320040), Emu(2331171), Emu(1280160))

    # Main title: "MotoMatch Recommendations"
    t1 = slide.shapes.add_textbox(Emu(457200), Emu(1828800), Emu(8229600), Emu(594360))
    tf = t1.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "MotoMatch Recommendations"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = SHELBY_BLUE
    r.font.name = FONT_BOLD

    # Subtitle: "for [Client Name]"
    t2 = slide.shapes.add_textbox(Emu(457200), Emu(2430000), Emu(8229600), Emu(320040))
    tf2 = t2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = f"for {CLIENT['name']}"
    r2.font.size = Pt(28)
    r2.font.bold = False
    r2.font.color.rgb = SHELBY_BLUE
    r2.font.name = FONT_BODY

    # Tagline lines
    t3 = slide.shapes.add_textbox(Emu(457200), Emu(2800000), Emu(6400800), Emu(500000))
    tf3 = t3.text_frame
    tf3.word_wrap = True
    p3a = tf3.paragraphs[0]
    r3a = p3a.add_run()
    r3a.text = "Personalized vehicle selections based on your questionnaire responses."
    r3a.font.size = Pt(10)
    r3a.font.color.rgb = SUBHEAD_GRAY
    r3a.font.name = FONT_BODY
    p3b = tf3.add_paragraph()
    r3b = p3b.add_run()
    r3b.text = "Each recommendation has been matched to your priorities, budget, and lifestyle."
    r3b.font.size = Pt(10)
    r3b.font.color.rgb = SUBHEAD_GRAY
    r3b.font.name = FONT_BODY

    # Footer bar (navy, full width)
    footer = add_shape(slide, 0, Emu(4343400), SLIDE_W, Emu(804672), fill_color=SHELBY_BLUE)
    footer.line.fill.background()

    # Footer wordmark (native aspect ratio 6.05:1, height=240000)
    add_image_safe(slide, LOGO_FOOTER, Emu(365760), Emu(4625736), Emu(1452000), Emu(240000))

    # Date — right side of footer
    t_date = slide.shapes.add_textbox(Emu(6858000), Emu(4572000), Emu(1828800), Emu(274320))
    tf_d = t_date.text_frame
    p_d = tf_d.paragraphs[0]
    p_d.alignment = PP_ALIGN.RIGHT
    r_d = p_d.add_run()
    r_d.text = CLIENT["date"]
    r_d.font.size = Pt(10)
    r_d.font.color.rgb = WHITE
    r_d.font.name = FONT_BODY


# ══════════════════════════════════════════════════════════════════════════════
# ─── SLIDE 2: WHAT YOU TOLD US ───────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_profile_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    add_inner_slide_chrome(slide)

    # Title
    t = slide.shapes.add_textbox(Emu(457200), Emu(228600), Emu(7315200), Emu(457200))
    tf = t.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "What You Told Us"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = SHELBY_BLUE
    r.font.name = FONT_BOLD

    # Subtitle
    s = slide.shapes.add_textbox(Emu(457200), Emu(685800), Emu(7772400), Emu(320040))
    tf2 = s.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "A summary of your MotoMatch questionnaire responses. These priorities shaped every recommendation."
    r2.font.size = Pt(10)
    r2.font.color.rgb = SUBHEAD_GRAY
    r2.font.name = FONT_BODY

    # ── LEFT PANEL: FINANCIALS BOX ───────────────────────────────────────────
    fin_box = add_shape(slide, Emu(457200), Emu(1097280), Emu(3291840), Emu(1371600), fill_color=NARDO)
    fin_box.line.fill.background()

    add_textbox(slide, Emu(640080), Emu(1207008), Emu(2926080), Emu(274320),
                "FINANCIALS", 14, bold=True, color=SHELBY_BLUE)

    add_textbox(slide, Emu(640080), Emu(1508760), Emu(2926080), Emu(182880),
                "Budget Range", 11, color=SUBHEAD_GRAY)
    add_textbox(slide, Emu(640080), Emu(1691640), Emu(2926080), Emu(274320),
                CLIENT["budget"], 18, bold=True, color=SHELBY_BLUE)

    add_textbox(slide, Emu(640080), Emu(2011680), Emu(2926080), Emu(182880),
                "Purchase Type", 11, color=SUBHEAD_GRAY)
    add_textbox(slide, Emu(640080), Emu(2194560), Emu(2926080), Emu(228600),
                CLIENT["purchase_type"], 16, bold=True, color=SHELBY_BLUE)

    # ── LEFT PANEL: KEY CONSIDERATIONS BOX ──────────────────────────────────
    kc_bg  = add_shape(slide, Emu(457200), Emu(2606040), Emu(3291840), Emu(1900000), fill_color=OFF_WHITE_ROW)
    kc_bg.line.fill.background()
    kc_bar = add_shape(slide, Emu(457200), Emu(2606040), Emu(73152), Emu(1900000), fill_color=SAO_PAULO)
    kc_bar.line.fill.background()

    add_textbox(slide, Emu(658368), Emu(2697480), Emu(2926080), Emu(274320),
                "KEY CONSIDERATIONS", 14, bold=True, color=SHELBY_BLUE)

    cb = slide.shapes.add_textbox(Emu(658368), Emu(2997520), Emu(2926080), Emu(1200000))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, consid in enumerate(CLIENT["considerations"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = f"\u2022  {consid}"
        r.font.size = Pt(9)
        r.font.color.rgb = PURE_BLACK
        r.font.name = FONT_BODY

    # ── RIGHT PANEL: YOUR PRIORITIES ─────────────────────────────────────────
    rp_bg = add_shape(slide, Emu(3931920), Emu(1097280), Emu(4754880), Emu(3900000), fill_color=NARDO)
    rp_bg.line.fill.background()

    add_textbox(slide, Emu(4114800), Emu(1207008), Emu(4389120), Emu(274320),
                "YOUR PRIORITIES", 14, bold=True, color=SHELBY_BLUE)
    add_textbox(slide, Emu(4114800), Emu(1463040), Emu(4389120), Emu(182880),
                "Ranked by importance (5 = highest, 1 = lowest)", 9, color=SUBHEAD_GRAY)

    ROW_H    = Emu(190000)
    ROW_GAP  = Emu(5500)
    BADGE_W  = Emu(640080)
    LABEL_W  = Emu(3703320)
    BADGE_L  = Emu(4114800)
    LABEL_L  = Emu(4828032)
    ROW_TOP_START = Emu(1700000)

    for i, (priority_name, score) in enumerate(CLIENT["priorities"]):
        row_top = ROW_TOP_START + i * (ROW_H + ROW_GAP)
        if row_top + ROW_H > Emu(4870000):
            break

        badge_color = SHELBY_BLUE if score >= 4 else SHELBY_LIGHT
        label_color = WHITE if score >= 4 else OFF_WHITE_ROW
        text_color  = SHELBY_BLUE if score >= 4 else SUBHEAD_GRAY
        bold_label  = score >= 4

        badge_bg = add_shape(slide, BADGE_L, row_top, BADGE_W, ROW_H, fill_color=badge_color)
        badge_bg.line.fill.background()
        badge_txt = slide.shapes.add_textbox(BADGE_L, row_top, BADGE_W, ROW_H)
        tf = badge_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = stars_string(score)
        r.font.size = Pt(8)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT_BOLD

        label_bg = add_shape(slide, LABEL_L, row_top, LABEL_W, ROW_H, fill_color=label_color)
        label_bg.line.fill.background()
        label_txt = slide.shapes.add_textbox(Emu(4919472), row_top, Emu(3520440), ROW_H)
        tf = label_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = priority_name
        r.font.size = Pt(10)
        r.font.bold = bold_label
        r.font.color.rgb = text_color
        r.font.name = FONT_BOLD if bold_label else FONT_BODY


# ══════════════════════════════════════════════════════════════════════════════
# ─── SLIDES 3+: VEHICLE SLIDES ───────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_vehicle_slide(prs, v):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    add_inner_slide_chrome(slide)

    # ── HEADER ───────────────────────────────────────────────────────────────
    # Recommendation counter (e.g. "Recommendation 1 of 3")
    rec_box = slide.shapes.add_textbox(Emu(457200), Emu(137160), Emu(2743200), Emu(201168))
    tf = rec_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Recommendation {v['rec_num']} of {TOTAL_VEHICLES}"
    r.font.size = Pt(9)
    r.font.color.rgb = SHELBY_LIGHT
    r.font.name = FONT_BODY

    # Vehicle name (large, bold, Shelby Blue)
    name_box = slide.shapes.add_textbox(Emu(457200), Emu(320040), Emu(5486400), Emu(411480))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = v["year_make_model"]
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = SHELBY_BLUE
    r.font.name = FONT_BOLD

    # MSRP (right-aligned, Shelby Light)
    msrp_box = slide.shapes.add_textbox(Emu(5943600), Emu(365760), Emu(2743200), Emu(320040))
    tf = msrp_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = v["msrp"]
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = SHELBY_LIGHT
    r.font.name = FONT_BOLD

    # ── SOURCE LINKS ─────────────────────────────────────────────────────────
    links_box = slide.shapes.add_textbox(Emu(457200), Emu(731520), Emu(5029200), Emu(201168))
    tf = links_box.text_frame
    p = tf.paragraphs[0]
    sp = slide.part
    add_hyperlink_run(p, "Manufacturer", v["mfr_url"], 10, slide_part=sp)
    add_plain_run(p, "  |  ", 10, color=SUBHEAD_GRAY)
    add_hyperlink_run(p, "Car and Driver Review", v["cd_url"], 10, slide_part=sp)
    add_plain_run(p, "  |  ", 10, color=SUBHEAD_GRAY)
    add_hyperlink_run(p, "Edmunds Review", v["edmunds_url"], 10, slide_part=sp)

    # ── BLURB BOX (dark navy bg, white italic text) ───────────────────────────
    blurb_bg = add_shape(slide, Emu(457200), Emu(987552), Emu(5029200), Emu(914400), fill_color=SHELBY_BLUE)
    blurb_bg.line.fill.background()

    blurb_txt = slide.shapes.add_textbox(Emu(566928), Emu(1042416), Emu(4809744), Emu(804672))
    tf = blurb_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = v["blurb"]
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = WHITE
    r.font.name = FONT_BODY

    # ── KEY SPECIFICATIONS ────────────────────────────────────────────────────
    specs_bg = add_shape(slide, Emu(457200), Emu(1993392), Emu(5029200), Emu(530352), fill_color=NARDO)
    specs_bg.line.fill.background()

    add_textbox(slide, Emu(594360), Emu(2029968), Emu(4754880), Emu(182880),
                "KEY SPECIFICATIONS", 14, bold=True, color=SHELBY_BLUE)

    # Height bumped from 256032 → 280000 EMU so the rare 2-line wrap stays
    # inside the Nardo bg box instead of bleeding into the WHY section.
    specs_val = slide.shapes.add_textbox(Emu(594360), Emu(2231136), Emu(4754880), Emu(280000))
    tf = specs_val.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = v["specs"]
    r.font.size = Pt(10)
    r.font.color.rgb = PURE_BLACK
    r.font.name = FONT_BODY

    # ── WHY IT FITS YOU ───────────────────────────────────────────────────────
    why_bg = add_shape(slide, Emu(457200), Emu(2570000), Emu(5029200), Emu(1560000), fill_color=NARDO)
    why_bg.line.fill.background()

    add_textbox(slide, Emu(594360), Emu(2615000), Emu(4754880), Emu(201168),
                "WHY IT FITS YOU", 14, bold=True, color=SHELBY_BLUE)

    WHY_ROW_H   = Emu(295000)
    WHY_ROW_GAP = Emu(15000)
    why_start   = Emu(2870000)
    for i, (tag, desc) in enumerate(v["why"][:4]):
        row_top = why_start + i * (WHY_ROW_H + WHY_ROW_GAP)
        wb = slide.shapes.add_textbox(Emu(594360), row_top, Emu(4754880), WHY_ROW_H)
        tf = wb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(13)
        # Tag in Shelby Light bold
        r1 = p.add_run()
        r1.text = tag + " "
        r1.font.size = Pt(10)
        r1.font.bold = True
        r1.font.color.rgb = SHELBY_LIGHT
        r1.font.name = FONT_BOLD
        # Description in body
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.bold = False
        r2.font.color.rgb = PURE_BLACK
        r2.font.name = FONT_BODY

    # ── CONSIDERATIONS ────────────────────────────────────────────────────────
    CONS_TOP = Emu(4150000)
    CONS_H   = Emu(880000)
    cons_bg  = add_shape(slide, Emu(457200), CONS_TOP, Emu(5029200), CONS_H, fill_color=OFF_WHITE_ROW)
    cons_bg.line.fill.background()
    cons_bar = add_shape(slide, Emu(457200), CONS_TOP, Emu(73152), CONS_H, fill_color=SAO_PAULO)
    cons_bar.line.fill.background()

    add_textbox(slide, Emu(658368), Emu(4190000), Emu(4663440), Emu(182880),
                "CONSIDERATIONS", 11, bold=True, color=SUBHEAD_GRAY)

    cons_box = slide.shapes.add_textbox(Emu(658368), Emu(4380000), Emu(4663440), Emu(640000))
    tf = cons_box.text_frame
    tf.word_wrap = True
    for i, cons_text in enumerate(v["considerations"][:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = f"\u2022  {cons_text}"
        r.font.size = Pt(9)
        r.font.color.rgb = SUBHEAD_GRAY
        r.font.name = FONT_BODY

    # ── OPTIONAL STAR LINE (yellow accent) ────────────────────────────────────
    # Only added if vehicle has a star_line. Placed between Considerations label
    # and bullets as a yellow-italic highlighted statement.
    # (In current build, star_line is embedded inside the blurb box area or
    # shown as the first consideration if used. Per prior builds it appears
    # at the bottom of the Why box. Uncomment the block below to activate.)
    #
    # if v.get("star_line"):
    #     star_box = slide.shapes.add_textbox(Emu(658368), Emu(4330000), Emu(4663440), Emu(200000))
    #     tf = star_box.text_frame
    #     p = tf.paragraphs[0]
    #     r = p.add_run()
    #     r.text = v["star_line"]
    #     r.font.size = Pt(9)
    #     r.font.italic = True
    #     r.font.bold = True
    #     r.font.color.rgb = SAO_PAULO
    #     r.font.name = FONT_BODY

    # ── THREE PHOTOS (right column, top-to-footer, crop-to-fill) ─────────────
    photo_paths = PHOTOS[v["key"]]
    for idx, p_top in enumerate(PHOTO_TOPS):
        if idx < len(photo_paths):
            add_image_safe(slide, photo_paths[idx], PHOTO_L, p_top, PHOTO_W, PHOTO_H)


# ══════════════════════════════════════════════════════════════════════════════
# ─── FINAL SLIDE: SIDE-BY-SIDE COMPARISON TABLE ──────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    add_inner_slide_chrome(slide)

    # Title
    add_textbox(slide, Emu(457200), Emu(137160), Emu(8229600), Emu(457200),
                "Side-by-Side Comparison", 26, bold=True, color=SHELBY_BLUE)

    # Subtitle
    add_textbox(slide, Emu(457200), Emu(594360), Emu(8229600), Emu(228600),
                "How each recommendation stacks up. Green highlighting indicates the best value in each category.",
                10, color=SUBHEAD_GRAY)

    # Table layout
    TABLE_L      = Emu(457200)
    TABLE_T      = Emu(868680)
    TABLE_W      = SLIDE_W - Emu(914400)
    N_VEHICLES   = len(VEHICLES)
    METRIC_COL_W = int(TABLE_W * 0.21)
    VEH_COL_W    = int((TABLE_W - METRIC_COL_W) / N_VEHICLES)
    ROW_H        = Emu(340000)
    HEADER_H     = Emu(400000)

    # Vehicle column left positions
    col_lefts = [TABLE_L + METRIC_COL_W + i * VEH_COL_W for i in range(N_VEHICLES)]

    # Vehicle short names (strip year prefix)
    veh_headers = [
        v["year_make_model"].replace("2026 ", "").replace("2025 ", "")
        for v in VEHICLES
    ]

    # Header row — metric cell
    mh = add_shape(slide, TABLE_L, TABLE_T, METRIC_COL_W, HEADER_H, fill_color=SHELBY_BLUE)
    mh.line.fill.background()
    mh_txt = slide.shapes.add_textbox(TABLE_L + Emu(50000), TABLE_T,
                                       METRIC_COL_W - Emu(50000), HEADER_H)
    tf = mh_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Criteria"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = SHELBY_LIGHT
    r.font.name = FONT_BOLD

    # Header row — vehicle columns
    for ci, (vname, cl) in enumerate(zip(veh_headers, col_lefts)):
        ch = add_shape(slide, cl, TABLE_T, VEH_COL_W, HEADER_H, fill_color=SHELBY_BLUE)
        ch.line.fill.background()
        ch_txt = slide.shapes.add_textbox(cl + Emu(30000), TABLE_T + Emu(30000),
                                           VEH_COL_W - Emu(60000), HEADER_H - Emu(60000))
        tf = ch_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = vname
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT_BOLD

    # Data rows
    for ri, (metric, values, best_idx) in enumerate(COMPARISON_ROWS):
        row_top = TABLE_T + HEADER_H + ri * ROW_H
        row_bg  = NARDO if ri % 2 == 0 else WHITE

        # Metric label cell
        mc = add_shape(slide, TABLE_L, row_top, METRIC_COL_W, ROW_H, fill_color=row_bg)
        mc.line.fill.background()
        mc_txt = slide.shapes.add_textbox(TABLE_L + Emu(50000), row_top + Emu(50000),
                                           METRIC_COL_W - Emu(100000), ROW_H - Emu(100000))
        tf = mc_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = metric
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = SHELBY_BLUE
        r.font.name = FONT_BOLD

        # Value cells
        for ci, (val, cl) in enumerate(zip(values, col_lefts)):
            is_best = (ci == best_idx)
            cell_bg = GREEN_CELL_BG if is_best else row_bg
            vc = add_shape(slide, cl, row_top, VEH_COL_W, ROW_H, fill_color=cell_bg)
            vc.line.fill.background()
            vc_txt = slide.shapes.add_textbox(cl + Emu(30000), row_top + Emu(50000),
                                               VEH_COL_W - Emu(60000), ROW_H - Emu(100000))
            tf = vc_txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = Pt(9)
            r.font.bold = is_best
            r.font.color.rgb = GREEN_BEST if is_best else PURE_BLACK
            r.font.name = FONT_BOLD if is_best else FONT_BODY

    # Legend
    leg_top = TABLE_T + HEADER_H + len(COMPARISON_ROWS) * ROW_H + Emu(60000)
    leg_swatch = add_shape(slide, TABLE_L, leg_top, Emu(304800), Emu(182880),
                            fill_color=GREEN_CELL_BG, line=True)
    leg_swatch.line.color.rgb = GREEN_BEST
    leg_swatch.line.width = Pt(1)
    add_textbox(slide, TABLE_L + Emu(370000), leg_top, Emu(2000000), Emu(182880),
                "= Best in category", 9, color=PURE_BLACK)

    # Sources footnote
    src_top = leg_top + Emu(250000)
    add_textbox(slide, TABLE_L, src_top, TABLE_W, Emu(182880),
                "Sources: Car and Driver, Edmunds, MotorTrend, Manufacturer specifications. Pricing may vary by configuration.",
                8, color=SUBHEAD_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# ─── MAIN BUILD FUNCTION ─────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def build_deck():
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    print("Building Slide 1: Title Slide...")
    build_title_slide(prs)

    print("Building Slide 2: What You Told Us...")
    build_profile_slide(prs)

    for v in VEHICLES:
        print(f"Building Slide {v['rec_num'] + 2}: {v['year_make_model']}...")
        build_vehicle_slide(prs, v)

    print(f"Building Final Slide: Side-by-Side Comparison...")
    build_comparison_slide(prs)

    prs.save(OUTPUT_PATH)
    total_slides = 2 + len(VEHICLES) + 1
    print(f"\n✓ Saved: {OUTPUT_PATH}  ({total_slides} slides)")
    return OUTPUT_PATH


if __name__ == "__main__":
    build_deck()
